import pdfminer.high_level as pdfh
import docx
import pptx
import sys

def ppt2txt(file_path):
    prs = pptx.Presentation(sys.argv[1])
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return '\n'.join(text_runs)
 
if __name__ == "__main__":
    try:
        result = ''
        if sys.argv[1][-3:] == 'pdf':
            print(sys.argv[1][-3:])
            result = pdfh.extract_text(sys.argv[1])
        elif sys.argv[1][-4:] == 'docx' or sys.argv[1][-3:] == 'doc':
            doc = docx.Document(sys.argv[1])
            result = [p.text for p in doc.paragraphs]
        elif sys.argv[1][-4:] == 'pptx' or sys.argv[1][-3:] == 'ppt':
            result = ppt2txt(sys.argv[1])
        print(result)

    except:
        print("ERR NO FILE, FORMAT: python DocScrapper.py <pdf U doc>")
        